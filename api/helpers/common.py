
from configparser import ConfigParser
import os
import pdfplumber
from docx import Document
from pptx import Presentation


def can_convert_int(s):
    try: 
        i = int(s)
        return True
    except:
        return False



def extract_text_from_pdf(file_path):
    text = ""
    with pdfplumber.open(file_path) as pdf:
        for page in pdf.pages:
            text += page.extract_text() + "\n"
    return text


def extract_text_from_docx(file_path):
    doc = Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text


def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text



def get_env_var(group, var_name): 
    config = ConfigParser()
    file_path = ".env"
    if os.path.isfile(file_path):
        config.read(file_path)
        return config[group][var_name]
    return os.environ.get(var_name)